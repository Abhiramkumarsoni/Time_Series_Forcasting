"""
generate_ppt.py

Generates a professional PowerPoint presentation for the
Retail Sales Forecasting project — ready for interviews.

Run:
    python generate_ppt.py
Output:
    Retail_Sales_Forecasting_Presentation.pptx  (in project root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette ─────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BG_CARD       = RGBColor(0x11, 0x25, 0x3B)   # card navy
ACCENT_BLUE   = RGBColor(0x00, 0x8B, 0xD8)   # electric blue
ACCENT_CYAN   = RGBColor(0x00, 0xD4, 0xFF)   # cyan
ACCENT_GREEN  = RGBColor(0x00, 0xE5, 0x96)   # mint green
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)   # amber
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)   # purple
TEXT_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT    = RGBColor(0xB0, 0xC4, 0xDE)   # light steel blue
TEXT_GREY     = RGBColor(0x78, 0x90, 0xA8)


# ── Helper utilities ───────────────────────────────────────────────────────────

def set_bg(slide, prs, color: RGBColor):
    """Fill the entire slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=TEXT_WHITE, align=PP_ALIGN.LEFT, italic=False):
    """Add a textbox with full styling."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines: list,
                           font_size=12, color=TEXT_WHITE):
    """Add a textbox with multiple bullet lines."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def accent_bar(slide, color: RGBColor, width=10.0, height=0.05, top=0.82):
    """Add a thin horizontal accent bar near the top."""
    add_rect(slide, 0, top, width, height, color)


def slide_title(slide, title: str, subtitle: str = "", accent=ACCENT_BLUE):
    """Draw the standard slide title block."""
    accent_bar(slide, accent)
    add_textbox(slide, 0.4, 0.1, 9.0, 0.6, title,
                font_size=28, bold=True, color=TEXT_WHITE)
    if subtitle:
        add_textbox(slide, 0.4, 0.65, 9.0, 0.35, subtitle,
                    font_size=13, color=TEXT_LIGHT, italic=True)


def card(slide, left, top, width, height, title, body_lines,
         accent_color=ACCENT_BLUE, title_size=14, body_size=11):
    """Draw a card with a coloured top border, title, and bullet list."""
    add_rect(slide, left, top, width, 0.06, accent_color)          # top strip
    add_rect(slide, left, top + 0.06, width, height - 0.06, BG_CARD)
    add_textbox(slide, left + 0.15, top + 0.1, width - 0.3, 0.35,
                title, font_size=title_size, bold=True, color=accent_color)
    add_multiline_textbox(slide, left + 0.15, top + 0.45,
                          width - 0.3, height - 0.55,
                          body_lines, font_size=body_size, color=TEXT_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)   # 16:9 widescreen

blank_layout = prs.slide_layouts[6]   # completely blank layout


# ─── SLIDE 1 — TITLE ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, prs, BG_DARK)

# Decorative gradient bar (simulated with stacked thin rects)
colors_gradient = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN]
bar_w = 10 / 3
for i, c in enumerate(colors_gradient):
    add_rect(s1, i * bar_w, 0, bar_w, 0.12, c)

# Main title
add_textbox(s1, 0.5, 0.9, 9.0, 1.0,
            "Retail Sales Forecasting",
            font_size=42, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# Sub-title
add_textbox(s1, 0.5, 1.85, 9.0, 0.5,
            "End-to-End Time Series ML System",
            font_size=20, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# Divider
add_rect(s1, 3.5, 2.45, 3.0, 0.04, ACCENT_BLUE)

# Tagline
add_textbox(s1, 0.5, 2.65, 9.0, 0.4,
            "ARIMA  ·  Prophet  ·  XGBoost  ·  LSTM  ·  Streamlit Dashboard",
            font_size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Bottom badge
add_rect(s1, 3.2, 4.5, 3.6, 0.55, ACCENT_BLUE)
add_textbox(s1, 3.2, 4.53, 3.6, 0.5,
            "Data Science  |  Statistical Modelling  |  Forecasting",
            font_size=10, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


# ─── SLIDE 2 — PROBLEM STATEMENT ──────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, prs, BG_DARK)
slide_title(s2, "Problem Statement", "Why does accurate forecasting matter in retail?", ACCENT_BLUE)

# Two big cards
card(s2, 0.3, 1.0, 4.4, 3.2,
     "📉  Under-Forecast",
     ["• Stockouts — shelves run empty",
      "• Lost revenue & lost customers",
      "• Poor service-level agreements",
      "• Emergency replenishment at premium cost",
      "• Brand reputation damage"],
     accent_color=RGBColor(0xE7, 0x4C, 0x3C))

card(s2, 5.3, 1.0, 4.4, 3.2,
     "📈  Over-Forecast",
     ["• Excess inventory ties up working capital",
      "• Holding cost: 20–30% of value per year",
      "• Warehouse space constraints",
      "• Forced markdowns erode margins",
      "• Waste for perishable goods"],
     accent_color=ACCENT_ORANGE)

# Centre solution arrow
add_textbox(s2, 4.0, 2.2, 2.0, 0.7,
            "✅\nOptimal\nForecast",
            font_size=11, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)


# ─── SLIDE 3 — PROJECT ARCHITECTURE ───────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, prs, BG_DARK)
slide_title(s3, "Project Architecture", "End-to-end modular ML pipeline", ACCENT_CYAN)

stages = [
    ("📥\nData\nLoader",     ACCENT_BLUE),
    ("🔧\nPre-\nprocessing", ACCENT_CYAN),
    ("⚙️\nFeature\nEngg.",    ACCENT_GREEN),
    ("🤖\nModel\nTraining",  ACCENT_ORANGE),
    ("📊\nStream-\nlit App",  ACCENT_PURPLE),
]
box_w, box_h = 1.55, 1.6
start_x = 0.3
for i, (label, color) in enumerate(stages):
    x = start_x + i * (box_w + 0.25)
    add_rect(s3, x, 1.3, box_w, box_h, color)
    add_textbox(s3, x, 1.3, box_w, box_h, label,
                font_size=13, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_textbox(s3, x + box_w, 1.85, 0.28, 0.5, "→",
                    font_size=20, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# File mapping row
file_labels = [
    "data_loader.py",
    "preprocessing.py",
    "feature_engineering.py",
    "trainer.py",
    "streamlit_app/",
]
for i, fname in enumerate(file_labels):
    x = start_x + i * (box_w + 0.25)
    add_textbox(s3, x, 3.0, box_w, 0.4, fname,
                font_size=9, color=TEXT_GREY, align=PP_ALIGN.CENTER, italic=True)

# Bottom note
add_textbox(s3, 0.3, 4.8, 9.4, 0.35,
            "All modules follow Single Responsibility Principle · Shared config.py for all hyperparameters",
            font_size=10, color=TEXT_GREY, align=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 4 — DATA & PREPROCESSING ──────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, prs, BG_DARK)
slide_title(s4, "Data & Preprocessing Pipeline", "4 years · 1,461 daily observations · No external download needed", ACCENT_GREEN)

card(s4, 0.3, 1.0, 3.0, 3.3,
     "📊  Synthetic Data",
     ["• 4 years of daily sales",
      "• Upward trend (+0.05/day)",
      "• Weekend uplift (+25 units)",
      "• Holiday season boost (+40)",
      "• Gaussian noise",
      "• Random promo spikes"],
     accent_color=ACCENT_GREEN, body_size=10)

card(s4, 3.5, 1.0, 3.0, 3.3,
     "🔧  Preprocessing Steps",
     ["• Full date range reindex",
      "• Time-based interpolation",
      "• IQR × 3 outlier clipping",
      "• Chronological 80/20 split",
      "• MinMaxScaler (fit on train only)",
      "• ⚠️  No data leakage!"],
     accent_color=ACCENT_CYAN, body_size=10)

card(s4, 6.7, 1.0, 3.0, 3.3,
     "🔑  Why Chron. Split?",
     ["• Random split = data leakage",
      "• Test set must be in future",
      "• Scaler fit on train only",
      "• Simulates real deployment:",
      "  train on past, predict future",
      "• Same logic as production"],
     accent_color=ACCENT_ORANGE, body_size=10)


# ─── SLIDE 5 — FEATURE ENGINEERING ───────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, prs, BG_DARK)
slide_title(s5, "Feature Engineering", "Creating temporal context for tree-based and ML models", ACCENT_ORANGE)

# Feature table header
headers = ["Feature", "Code", "Business Meaning"]
col_x   = [0.3, 3.5, 6.7]
col_w   = [3.0, 3.0, 3.0]

add_rect(s5, 0.3, 1.0, 9.4, 0.4, ACCENT_ORANGE)
for i, h in enumerate(headers):
    add_textbox(s5, col_x[i] + 0.1, 1.0, col_w[i] - 0.1, 0.4,
                h, font_size=12, bold=True, color=TEXT_WHITE)

rows = [
    ("lag_1",              "Sales.shift(1)",          "Yesterday's sales — momentum"),
    ("lag_7",              "Sales.shift(7)",           "Same weekday last week"),
    ("lag_14 / lag_30",    "Sales.shift(14/30)",      "Bi-weekly & monthly patterns"),
    ("rolling_mean_7",     ".shift(1).rolling(7).mean()", "7-day smoothed average"),
    ("rolling_std_14",     ".shift(1).rolling(14).std()","Volatility / uncertainty"),
    ("is_weekend",         "dayofweek >= 5",           "Sat/Sun sales uplift"),
    ("is_holiday_season",  "month in [11, 12]",        "Black Friday / Christmas"),
]

row_colors = [BG_CARD, BG_DARK]
for r, (feat, code, meaning) in enumerate(rows):
    y = 1.45 + r * 0.42
    add_rect(s5, 0.3, y, 9.4, 0.42, row_colors[r % 2])
    add_textbox(s5, col_x[0] + 0.1, y + 0.05, 2.8, 0.35, feat,  font_size=10, bold=True,  color=ACCENT_CYAN)
    add_textbox(s5, col_x[1] + 0.1, y + 0.05, 2.8, 0.35, code,  font_size=9,  italic=True, color=TEXT_LIGHT)
    add_textbox(s5, col_x[2] + 0.1, y + 0.05, 2.8, 0.35, meaning, font_size=10, color=TEXT_WHITE)

# Key note
add_textbox(s5, 0.3, 4.5, 9.4, 0.4,
            "⚠️  .shift(1) before rolling prevents look-ahead leakage — future values never touch the model during training",
            font_size=10, color=ACCENT_ORANGE, italic=True, align=PP_ALIGN.CENTER)


# ─── SLIDE 6 — MODELS OVERVIEW ────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, prs, BG_DARK)
slide_title(s6, "Four Models — Four Paradigms", "Why compare? Each model solves a different aspect of the problem", ACCENT_PURPLE)

model_cards = [
    ("📐  ARIMA (5,1,2)",
     ["Statistical baseline",
      "Interpretable & fast",
      "Needs stationary series",
      "Best: short data, explainability",
      "Output: confidence intervals"],
     ACCENT_BLUE),
    ("🔮  Prophet",
     ["Additive decomposition",
      "Handles seasonality auto",
      "Built-in holiday effects",
      "Best: business reporting",
      "Output: trend + seasonal"],
     ACCENT_CYAN),
    ("🌲  XGBoost",
     ["Tree-based, row-independent",
      "Needs lag features explicitly",
      "Fast & highly accurate",
      "Best: large datasets",
      "Output: point forecast"],
     ACCENT_GREEN),
    ("🧠  LSTM",
     ["Deep learning, sequential",
      "Look-back window = 30 days",
      "Learns non-linear patterns",
      "Best: complex long series",
      "Output: sequence prediction"],
     ACCENT_ORANGE),
]

for i, (title, bullets, color) in enumerate(model_cards):
    x = 0.3 + i * 2.42
    card(s6, x, 1.0, 2.25, 3.8, title, bullets,
         accent_color=color, title_size=12, body_size=9)

# Shared interface note
add_rect(s6, 0.3, 5.05, 9.4, 0.45, BG_CARD)
add_textbox(s6, 0.3, 5.07, 9.4, 0.4,
            "✅  All models share the same interface:  fit()  ·  predict()  ·  save()  ·  load()   → Strategy Design Pattern",
            font_size=11, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)


# ─── SLIDE 7 — ARIMA DEEP DIVE ────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank_layout)
set_bg(s7, prs, BG_DARK)
slide_title(s7, "ARIMA Deep Dive  —  ARIMA(5, 1, 2)", "The mathematical foundation of classical time series forecasting", ACCENT_BLUE)

# Three parameter boxes
params = [
    ("p = 5\nAuto-Regressive",
     "Today depends on the last 5 days.\nThe model learns weights for Y(t-1) … Y(t-5)",
     ACCENT_BLUE),
    ("d = 1\nIntegrated",
     "Apply first-order differencing.\nY'(t) = Y(t) - Y(t-1)  →  removes the trend.\nVerified via ADF Test (p-value < 0.05)",
     ACCENT_CYAN),
    ("q = 2\nMoving Average",
     "Uses past 2 forecast errors\nε(t-1) and ε(t-2) to correct predictions.\nMakes model self-correcting.",
     ACCENT_GREEN),
]

for i, (heading, body, color) in enumerate(params):
    x = 0.3 + i * 3.25
    add_rect(s7, x, 1.0, 3.0, 0.7, color)
    add_textbox(s7, x, 1.0, 3.0, 0.7, heading,
                font_size=15, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s7, x, 1.7, 3.0, 1.6, BG_CARD)
    add_textbox(s7, x + 0.1, 1.75, 2.8, 1.5, body,
                font_size=10, color=TEXT_LIGHT)

# Diagnostics row
add_textbox(s7, 0.3, 3.5, 9.4, 0.35,
            "📋  Model Diagnostics — How We Know It Works",
            font_size=13, bold=True, color=TEXT_WHITE)

diag_items = [
    ("ADF Test", "p-value < 0.05 → Series is stationary after d=1 differencing", ACCENT_BLUE),
    ("Ljung-Box Test", "p-value > 0.05 → Residuals are white noise (no pattern left)", ACCENT_CYAN),
    ("ACF / PACF Plots", "PACF cuts at lag 5 → p=5.  ACF cuts at lag 2 → q=2", ACCENT_GREEN),
]

for i, (test, desc, color) in enumerate(diag_items):
    x = 0.3 + i * 3.25
    add_rect(s7, x, 3.95, 3.0, 0.05, color)
    add_rect(s7, x, 4.0, 3.0, 1.1, BG_CARD)
    add_textbox(s7, x + 0.1, 4.05, 2.8, 0.3, test, font_size=11, bold=True, color=color)
    add_textbox(s7, x + 0.1, 4.38, 2.8, 0.7, desc, font_size=10, color=TEXT_LIGHT)


# ─── SLIDE 8 — SARIMA SEASONAL ────────────────────────────────────────────────
s8 = prs.slides.add_slide(blank_layout)
set_bg(s8, prs, BG_DARK)
slide_title(s8, "SARIMA Seasonal Order  —  (1, 1, 1, 7)", "Adding weekly seasonality to capture weekend uplift patterns", ACCENT_CYAN)

seasonal_params = [
    ("P = 1\nSeasonal AR",
     "Depends on sales from\nexactly 7 days ago (same weekday).\nYesterday's Monday predicts today's Monday.",
     ACCENT_BLUE),
    ("D = 1\nSeasonal Diff",
     "Seasonal differencing:\nY'(t) = Y(t) - Y(t-7)\nRemoves the repeating weekly pattern.",
     ACCENT_CYAN),
    ("Q = 1\nSeasonal MA",
     "Uses forecast error from\n7 days ago to self-correct\nweekly prediction bias.",
     ACCENT_GREEN),
    ("s = 7\nPeriod",
     "The seasonal cycle = 7 days.\nCaptures Mon–Sun weekly\nsales rhythm in retail.",
     ACCENT_ORANGE),
]

for i, (heading, body, color) in enumerate(seasonal_params):
    x = 0.25 + i * 2.42
    add_rect(s8, x, 1.0, 2.2, 0.7, color)
    add_textbox(s8, x, 1.0, 2.2, 0.7, heading,
                font_size=13, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s8, x, 1.7, 2.2, 1.6, BG_CARD)
    add_textbox(s8, x + 0.1, 1.8, 2.0, 1.4, body, font_size=10, color=TEXT_LIGHT)

# Why weekly makes sense box
add_rect(s8, 0.3, 3.55, 9.4, 1.65, BG_CARD)
add_rect(s8, 0.3, 3.55, 9.4, 0.05, ACCENT_CYAN)
add_textbox(s8, 0.5, 3.65, 9.0, 0.35,
            "📊  Why s=7 Makes Business Sense",
            font_size=12, bold=True, color=ACCENT_CYAN)
add_multiline_textbox(s8, 0.5, 4.05, 9.0, 1.0,
    ["• Retail sales exhibit strong weekly cycles — weekends (Sat/Sun) consistently outperform weekdays",
     "• SARIMA with s=7 directly models this rhythm without needing explicit 'is_weekend' feature (unlike XGBoost)",
     "• The seasonal differencing (D=1) removes the weekly pattern so the non-seasonal part (p,d,q) models the residual trend"],
    font_size=10, color=TEXT_LIGHT)


# ─── SLIDE 9 — EVALUATION METRICS ────────────────────────────────────────────
s9 = prs.slides.add_slide(blank_layout)
set_bg(s9, prs, BG_DARK)
slide_title(s9, "Model Evaluation Metrics", "How we measure forecast quality — and why each metric matters", ACCENT_GREEN)

metric_cards = [
    ("RMSE",
     "Root Mean Squared Error",
     "Penalises large errors heavily.\nUse for model selection & training.\nDirectly linked to safety stock calculation.",
     "σ_error ≈ RMSE",
     ACCENT_BLUE),
    ("MAE",
     "Mean Absolute Error",
     "Average absolute deviation.\nMore robust to outliers than RMSE.\nEasy to interpret in same units as sales.",
     "avg |actual - predicted|",
     ACCENT_CYAN),
    ("MAPE",
     "Mean Abs % Error",
     "Express error as a percentage.\nEasy for business stakeholders.\n'Our forecast is off by X% on average'",
     "∑|actual-pred|/actual × 100",
     ACCENT_GREEN),
    ("R²",
     "Coefficient of Determination",
     "% variance explained by model.\nR²=0.95 → model explains 95%.\nCan be misleading alone in TS.",
     "1 − SS_res/SS_tot",
     ACCENT_ORANGE),
]

for i, (name, full, body, formula, color) in enumerate(metric_cards):
    x = 0.3 + i * 2.42
    add_rect(s9, x, 1.0, 2.25, 0.55, color)
    add_textbox(s9, x, 1.0, 2.25, 0.55, name,
                font_size=22, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s9, x, 1.55, 2.25, 3.35, BG_CARD)
    add_textbox(s9, x + 0.1, 1.6, 2.05, 0.35, full, font_size=10, bold=True, color=color)
    add_textbox(s9, x + 0.1, 2.0, 2.05, 1.6, body, font_size=9.5, color=TEXT_LIGHT)
    add_rect(s9, x + 0.1, 3.65, 2.05, 0.05, color)
    add_textbox(s9, x + 0.1, 3.73, 2.05, 0.5, formula,
                font_size=9, color=color, italic=True, align=PP_ALIGN.CENTER)

# Business link
add_rect(s9, 0.3, 5.0, 9.4, 0.48, BG_CARD)
add_rect(s9, 0.3, 5.0, 9.4, 0.05, ACCENT_GREEN)
add_textbox(s9, 0.5, 5.07, 9.0, 0.38,
            "💡  Safety Stock = Z × √(Lead Time) × RMSE   →   Lower RMSE = less safety stock = freed working capital",
            font_size=11, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)


# ─── SLIDE 10 — STREAMLIT DASHBOARD ──────────────────────────────────────────
s10 = prs.slides.add_slide(blank_layout)
set_bg(s10, prs, BG_DARK)
slide_title(s10, "Live Interactive Dashboard", "Built with Streamlit — 4 pages, zero-config, runs locally", ACCENT_PURPLE)

pages = [
    ("🏠  Home",
     ["KPI cards: Mean, Min, Max, Date range",
      "Interactive date range filter",
      "Full 4-year time series chart"],
     ACCENT_BLUE),
    ("📊  EDA",
     ["4 tabs: Seasonal, Monthly, Decomposition, Heatmap",
      "Trend + seasonal + residual decomposition",
      "Correlation heatmap of all features"],
     ACCENT_CYAN),
    ("🔮  Forecast",
     ["Choose model + forecast horizon (slider)",
      "Trains model live, shows 95% CI band",
      "Download forecast as CSV"],
     ACCENT_GREEN),
    ("📈  Model Comparison",
     ["Trains all 4 models simultaneously",
      "Metrics table with best highlighted in green",
      "Overlay chart of all 4 forecasts"],
     ACCENT_ORANGE),
]

for i, (page, bullets, color) in enumerate(pages):
    x = 0.3 + i * 2.42
    card(s10, x, 1.0, 2.25, 3.85, page, bullets,
         accent_color=color, title_size=12, body_size=9.5)

# Launch command
add_rect(s10, 0.3, 5.0, 9.4, 0.48, BG_CARD)
add_textbox(s10, 0.5, 5.07, 9.0, 0.38,
            "▶  streamlit run streamlit_app/Home.py   →   Opens at  http://localhost:8501",
            font_size=12, bold=True, color=ACCENT_GREEN,
            align=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 11 — ENGINEERING EXCELLENCE ───────────────────────────────────────
s11 = prs.slides.add_slide(blank_layout)
set_bg(s11, prs, BG_DARK)
slide_title(s11, "Engineering Excellence", "Production-quality code — not just a notebook", ACCENT_BLUE)

eng_cards = [
    ("🧪  Unit Tests",
     ["pytest with 25+ test cases",
      "test_models.py: fit/predict/save/load",
      "test_metrics.py: RMSE, MAE, MAPE, R²",
      "conftest.py for shared fixtures"],
     ACCENT_CYAN),
    ("🐳  Docker",
     ["Dockerfile using python:3.11-slim",
      "docker-compose.yml for one-command launch",
      "Health check endpoint included",
      "docker-compose up --build"],
     ACCENT_BLUE),
    ("⚙️  CI/CD",
     ["GitHub Actions workflow on push to main",
      "Runs ruff (linting) + black (formatting)",
      "Runs full pytest suite with coverage",
      "Tests on Python 3.10 and 3.11"],
     ACCENT_GREEN),
    ("📁  Project Structure",
     ["src/ — 9 modular Python files",
      "src/models/ — 4 model classes + trainer",
      "config.py — single source of truth",
      "Makefile for one-command workflows"],
     ACCENT_ORANGE),
]

for i, (title, bullets, color) in enumerate(eng_cards):
    x = 0.3 + (i % 2) * 4.85
    y = 1.0 + (i // 2) * 2.35
    card(s11, x, y, 4.5, 2.2, title, bullets,
         accent_color=color, title_size=12, body_size=10)


# ─── SLIDE 12 — BUSINESS IMPACT ──────────────────────────────────────────────
s12 = prs.slides.add_slide(blank_layout)
set_bg(s12, prs, BG_DARK)
slide_title(s12, "Business Impact & Real-World Value", "Translating model metrics into rupees saved", ACCENT_GREEN)

add_textbox(s12, 0.3, 0.95, 9.4, 0.4,
            "Every 1-unit reduction in RMSE directly translates to reduced safety stock requirements:",
            font_size=12, color=TEXT_LIGHT, italic=True)

# Formula block
add_rect(s12, 0.5, 1.4, 9.0, 0.9, BG_CARD)
add_rect(s12, 0.5, 1.4, 9.0, 0.06, ACCENT_GREEN)
add_textbox(s12, 0.5, 1.5, 9.0, 0.7,
            "Safety Stock  =  Z  ×  √(Lead Time)  ×  σ_error    where  σ_error  ≈  RMSE of validation set",
            font_size=15, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

impact_items = [
    ("10–15%\nReduction", "in safety stock when switching\nfrom naive baseline to XGBoost/Prophet", ACCENT_BLUE),
    ("20–30%\nHolding Cost", "of inventory value spent annually\non warehousing — directly reduced", ACCENT_ORANGE),
    ("95%\nService Level", "maintained with lower stock\nusing Z=1.65 safety factor", ACCENT_GREEN),
    ("1 Dashboard\nfor All Teams", "Operations, Finance, Supply Chain\nuse the same forecast output", ACCENT_PURPLE),
]

for i, (metric, desc, color) in enumerate(impact_items):
    x = 0.3 + i * 2.42
    add_rect(s12, x, 2.55, 2.25, 0.9, color)
    add_textbox(s12, x, 2.55, 2.25, 0.9, metric,
                font_size=16, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s12, x, 3.45, 2.25, 1.2, BG_CARD)
    add_textbox(s12, x + 0.1, 3.5, 2.05, 1.1, desc, font_size=9.5, color=TEXT_LIGHT)

add_textbox(s12, 0.3, 4.8, 9.4, 0.35,
            "This project demonstrates that a data scientist must speak the language of business, not just code.",
            font_size=11, color=ACCENT_ORANGE, italic=True, align=PP_ALIGN.CENTER)


# ─── SLIDE 13 — FUTURE IMPROVEMENTS ─────────────────────────────────────────
s13 = prs.slides.add_slide(blank_layout)
set_bg(s13, prs, BG_DARK)
slide_title(s13, "Future Improvements & Roadmap", "Demonstrating forward thinking — always impress interviewers with this", ACCENT_ORANGE)

improvements = [
    ("🔍  Optuna Hyperparameter Tuning",
     ["Auto-search best ARIMA orders (p,d,q)",
      "XGBoost: n_estimators, learning_rate, max_depth",
      "LSTM: units, look_back, batch_size",
      "Uses Bayesian optimisation (faster than grid search)"],
     ACCENT_BLUE),
    ("📦  MLflow Experiment Tracking",
     ["Log every training run: params + metrics + model",
      "Compare runs in a visual UI",
      "Register and version production models",
      "One-click model rollback if performance degrades"],
     ACCENT_GREEN),
    ("🤝  Ensemble Model",
     ["Combine all 4 forecasts using weighted average",
      "Or train a meta-learner (stacking)",
      "Ensembles almost always outperform individual models",
      "Reduce variance across model failures"],
     ACCENT_ORANGE),
    ("⚡  Real-Time Streaming",
     ["Apache Kafka or AWS Kinesis for live data",
      "Online learning: update model as new sales arrive",
      "Alerts when distribution shifts (drift detection)",
      "FastAPI endpoint for serving predictions"],
     ACCENT_PURPLE),
]

for i, (title, bullets, color) in enumerate(improvements):
    x = 0.3 + (i % 2) * 4.85
    y = 1.0 + (i // 2) * 2.35
    card(s13, x, y, 4.5, 2.15, title, bullets,
         accent_color=color, title_size=11, body_size=9.5)


# ─── SLIDE 14 — CLOSING / Q&A ─────────────────────────────────────────────────
s14 = prs.slides.add_slide(blank_layout)
set_bg(s14, prs, BG_DARK)

colors_gradient = [ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN]
bar_w = 10 / 3
for i, c in enumerate(colors_gradient):
    add_rect(s14, i * bar_w, 5.5, bar_w, 0.125, c)

add_textbox(s14, 0.5, 0.7, 9.0, 1.0,
            "Thank You",
            font_size=52, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

add_rect(s14, 3.5, 1.75, 3.0, 0.05, ACCENT_BLUE)

add_textbox(s14, 0.5, 1.95, 9.0, 0.5,
            "Retail Sales Forecasting  |  End-to-End Time Series ML System",
            font_size=15, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

summary = [
    "✅  4-Year Daily Sales Dataset (1,461 data points)",
    "✅  Full Preprocessing Pipeline  (interpolation · outlier clipping · no leakage)",
    "✅  4 Models Compared  (ARIMA · Prophet · XGBoost · LSTM)",
    "✅  Interactive Streamlit Dashboard  (Home · EDA · Forecast · Model Comparison)",
    "✅  Production-Ready  (pytest · Docker · GitHub Actions CI)",
]
add_multiline_textbox(s14, 2.0, 2.6, 6.0, 2.5, summary, font_size=12, color=TEXT_LIGHT)

add_textbox(s14, 0.5, 4.7, 9.0, 0.4,
            "Questions Welcome  🙌",
            font_size=20, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
output_path = "Retail_Forecast_Presentation_v2.pptx"
prs.save(output_path)
print("\nDONE! Presentation saved -> " + output_path)
print("Slides: " + str(len(prs.slides)))
print("Open with Microsoft PowerPoint or Google Slides\n")

